import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import os
import statsmodels.api as sm
import urllib.request
import json
import io
from streamlit_oauth import OAuth2Component

# ==========================================
# PAGE CONFIGURATION & CUSTOM CSS
# ==========================================
st.set_page_config(
    page_title="Impact Analytics Dashboard",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.markdown("""
<style>
    [data-testid="stMetricValue"] {
        font-size: 2rem;
        font-weight: 700;
    }
    .stTabs [data-baseweb="tab-list"] {
        gap: 24px;
    }
    .stTabs [data-baseweb="tab"] {
        height: 50px;
        white-space: pre-wrap;
        background-color: transparent;
        border-radius: 4px 4px 0px 0px;
        padding-top: 10px;
        padding-bottom: 10px;
    }
</style>
""", unsafe_allow_html=True)


# ==========================================
# FIX 4 — MODULE-LEVEL COLOR CONSTANTS
# Single source of truth; both modules import from here.
# ==========================================
COLOR_MAP = {'Baseline': '#636EFA', 'Endline': '#00CC96'}
RISE_COLORS = {
    "Reviving":  "#f27c48",
    "Initiating": "#0094c9",
    "Shaping":   "#00964d",
    "Evolving":  "#ed1c2d",
}
# Longitudinal module uses the same palette — one dict, two names.
RISE_COLORS_LONG = RISE_COLORS
TIME_ORDER = ['AY24-25 Baseline', 'AY24-25 Endline', 'AY25-26 Baseline', 'AY25-26 Endline']
AY_ORDER   = ['AY24-25', 'AY25-26']


# ==========================================
# AUTHENTICATION GATEKEEPER
# ==========================================
try:
    CLIENT_ID     = st.secrets["GOOGLE_CLIENT_ID"]
    CLIENT_SECRET = st.secrets["GOOGLE_CLIENT_SECRET"]
except FileNotFoundError:
    st.error(
        "Missing `.streamlit/secrets.toml` file or Streamlit Cloud Secrets. "
        "Please ensure your Google Client ID and Secret are configured."
    )
    st.stop()

AUTHORIZE_URL     = "https://accounts.google.com/o/oauth2/v2/auth"
TOKEN_URL         = "https://oauth2.googleapis.com/token"
REVOKE_TOKEN_URL  = "https://oauth2.googleapis.com/revoke"

if "logged_in_email"  not in st.session_state:
    st.session_state["logged_in_email"]  = None
if "user_first_name"  not in st.session_state:
    st.session_state["user_first_name"]  = "User"

if not st.session_state["logged_in_email"]:
    col1, col2, col3 = st.columns(3)
    with col2:
        st.write("")
        st.write("")
        try:
            st.image("evidyaloka_logo.png", width=320)
        except Exception:
            st.empty()
        st.markdown(
            "<h2 style='text-align:center;color:#0094c9;'>Student Analytics Portal</h2>",
            unsafe_allow_html=True,
        )
        st.markdown(
            "<p style='text-align:center;'>Please sign in with your @evidyaloka.org email to access the dashboard.</p>",
            unsafe_allow_html=True,
        )
        st.markdown("---")
        oauth2 = OAuth2Component(
            CLIENT_ID, CLIENT_SECRET, AUTHORIZE_URL, TOKEN_URL, REVOKE_TOKEN_URL
        )
        result = oauth2.authorize_button(
            name="Sign in with Google",
            icon="https://upload.wikimedia.org/wikipedia/commons/5/53/Google_%22G%22_Logo.svg",
            redirect_uri="https://ev-assessments.streamlit.app",
            scope="openid email profile",
            key="google_login",
            use_container_width=True,
        )
        if result and "token" in result:
            id_token = result["token"]["id_token"]
            # FIX 1 — id_token list bug: was `id_token = id_token` (no-op).
            if isinstance(id_token, list):
                id_token = id_token[0]
            verify_url = f"https://oauth2.googleapis.com/tokeninfo?id_token={id_token}"
            try:
                with urllib.request.urlopen(verify_url) as response:
                    user_info = json.loads(response.read().decode())
                st.session_state["logged_in_email"] = user_info.get("email")
                st.session_state["user_first_name"] = user_info.get("given_name", "User")
                st.rerun()
            except Exception as e:
                st.error(f"Error verifying login with Google: {e}")
                st.stop()
    st.stop()


# ==========================================
# FIX 3 — SHARED FILTER SIDEBAR BUILDER
# Replaces the near-identical 60-line blocks that appeared in both
# the main dashboard and the longitudinal module.
# Returns filtered_df and the dict of selected values.
# key_prefix keeps widget keys unique across pages.
# ==========================================
def build_filter_sidebar(df: pd.DataFrame, key_prefix: str) -> tuple[pd.DataFrame, dict]:
    """
    Renders cascading State → Donor → Centre → Subject → Grade → Gender
    filters in the sidebar and returns (filtered_df, selections).
    """
    selections = {}
    st.sidebar.header("🎯 Global Filters")

    # State
    states = ["All"] + sorted(df["State"].dropna().astype(str).unique()) if "State" in df.columns else ["All"]
    sel_state = st.sidebar.selectbox("Select State", states, index=0, key=f"{key_prefix}_state")
    selections["state"] = sel_state
    dff = df[df["State"].astype(str) == sel_state].copy() if sel_state != "All" else df.copy()

    # Donor
    donors = ["All"] + sorted(dff["Donor"].dropna().astype(str).unique()) if "Donor" in dff.columns else ["All"]
    sel_donor = st.sidebar.selectbox("Select Donor", donors, index=0, key=f"{key_prefix}_donor")
    selections["donor"] = sel_donor
    if sel_donor != "All":
        dff = dff[dff["Donor"].astype(str) == sel_donor]

    # Centre
    centres = ["All"] + sorted(dff["Centre Name"].dropna().astype(str).unique()) if "Centre Name" in dff.columns else ["All"]
    sel_centre = st.sidebar.selectbox("Select Centre", centres, index=0, key=f"{key_prefix}_centre")
    selections["centre"] = sel_centre
    if sel_centre != "All":
        dff = dff[dff["Centre Name"].astype(str) == sel_centre]

    # Subject
    subjects = ["All"] + sorted(dff["Subject"].dropna().astype(str).unique()) if "Subject" in dff.columns else ["All"]
    sel_subject = st.sidebar.selectbox("Select Subject", subjects, index=0, key=f"{key_prefix}_subject")
    selections["subject"] = sel_subject
    if sel_subject != "All":
        dff = dff[dff["Subject"].astype(str) == sel_subject]

    # Grade (multi-select; longitudinal module filters by AY24-25 cohort grade)
    if key_prefix == "long":
        df_base_year = dff[dff["Academic Year"] == "AY24-25"] if "Academic Year" in dff.columns else dff
        grades = sorted(df_base_year["Grade"].dropna().astype(str).unique()) if "Grade" in df_base_year.columns else []
        sel_grades = st.sidebar.multiselect(
            "Select AY 24-25 Grade (Cohort Tracking)", options=grades, default=grades,
            key=f"{key_prefix}_grade",
            help="Select the student's grade in AY 24-25. The dashboard will automatically track them into their promoted grade for AY 25-26.",
        )
        if sel_grades:
            cohort_ids = df_base_year[df_base_year["Grade"].astype(str).isin(sel_grades)]["Student ID"].unique()
            dff = dff[dff["Student ID"].isin(cohort_ids)]
        else:
            dff = dff.iloc[0:0]
    else:
        grades = sorted(dff["Grade"].dropna().astype(str).unique()) if "Grade" in dff.columns else []
        sel_grades = st.sidebar.multiselect(
            "Select Grade(s)", options=grades, default=grades, key=f"{key_prefix}_grade"
        )
        if sel_grades:
            dff = dff[dff["Grade"].astype(str).isin(sel_grades)]
        else:
            dff = dff.iloc[0:0]
    selections["grades"] = sel_grades

    # Gender
    if "Gender" in dff.columns:
        valid = dff.dropna(subset=["Gender"])
        valid = valid[~valid["Gender"].astype(str).str.lower().isin(["nan", "none", "null", ""])]
        genders = sorted(valid["Gender"].astype(str).unique())
        if genders:
            sel_genders = st.sidebar.multiselect(
                "Select Gender(s)", options=genders, default=genders, key=f"{key_prefix}_gender"
            )
            dff = dff[dff["Gender"].astype(str).isin(sel_genders)]
            selections["genders"] = sel_genders
        else:
            selections["genders"] = []
    else:
        selections["genders"] = []

    return dff, selections


# ==========================================
# APP ROUTER / HOMEPAGE GATEKEEPER
# ==========================================
if "current_page" not in st.session_state:
    st.session_state["current_page"] = "home"

if st.session_state["current_page"] == "home":
    st.write("")
    st.write("")
    st.title(f"👋 Welcome, {st.session_state['user_first_name']}!")
    st.markdown(
        "<p style='color:gray;font-size:1.1em;'>Select an application below to continue.</p>",
        unsafe_allow_html=True,
    )
    st.markdown("---")
    col1, col2, col3, col4 = st.columns(4)
    with col1:
        st.markdown("<h1 style='text-align:center;font-size:4rem;'>📈</h1>", unsafe_allow_html=True)
        if st.button("Impact Analytics Dashboard", use_container_width=True):
            st.session_state["current_page"] = "dashboard"
            st.rerun()
    with col2:
        st.markdown("<h1 style='text-align:center;font-size:4rem;'>🏛️</h1>", unsafe_allow_html=True)
        if st.button("Longitudinal Analysis", use_container_width=True):
            st.session_state["current_page"] = "longitudinal"
            st.rerun()
    st.stop()


# ==========================================
# LONGITUDINAL ANALYSIS MODULE
# ==========================================
if st.session_state["current_page"] == "longitudinal":
    st.title("🏛️ Strategic Longitudinal Analysis")
    st.markdown(
        "<p style='color:gray;font-size:1.1em;'>Year-over-Year Trajectories, Equity Tracking, "
        "and Strategic Insights (AY 24-25 vs AY 25-26)</p>",
        unsafe_allow_html=True,
    )
    st.markdown("---")

    @st.cache_data
    def load_multi_year_data(file_24, file_25):
        def clean_sheet(df, year, period):
            if df.empty:
                return pd.DataFrame()
            cols = [
                'State', 'Centre Name', 'Donor', 'Subject', 'Grade',
                'Student ID', 'Gender', 'Obtained Marks', 'Rubrics', 'Category',
            ]
            df_clean = df[[c for c in cols if c in df.columns]].copy()
            if 'Rubrics' in df_clean.columns:
                df_clean.rename(columns={'Rubrics': 'Category'}, inplace=True)

            df_clean['Academic Year'] = year
            df_clean['Period']        = period
            df_clean['Timepoint']     = f"{year} {period}"
            df_clean['Obtained Marks'] = pd.to_numeric(df_clean['Obtained Marks'], errors='coerce')

            # Robust Student ID cleaning (handles float, scientific notation,
            # alphanumeric IDs, and invisible Unicode characters)
            raw_ids     = df_clean['Student ID'].copy()
            numeric_ids = pd.to_numeric(raw_ids, errors='coerce')
            is_numeric  = numeric_ids.notna()

            df_clean['Student ID'] = pd.NA

            df_clean.loc[is_numeric, 'Student ID'] = (
                numeric_ids[is_numeric]
                .astype('Int64')
                .astype(str)
                .str.strip()
            )
            df_clean.loc[~is_numeric, 'Student ID'] = (
                raw_ids[~is_numeric]
                .astype(str)
                .str.strip()
                .str.upper()
            )

            invalid = df_clean['Student ID'].astype(str).str.upper().isin(
                ['NA', 'NAN', 'NONE', 'NULL', '<NA>', '']
            )
            df_clean.loc[invalid, 'Student ID'] = pd.NA

            for col in ['State', 'Centre Name', 'Donor', 'Subject', 'Gender', 'Category']:
                if col in df_clean.columns:
                    df_clean[col] = df_clean[col].astype(str).str.strip()
            if 'Gender' in df_clean.columns:
                df_clean['Gender'] = df_clean['Gender'].str.title()
            if 'Grade' in df_clean.columns:
                df_clean['Grade'] = df_clean['Grade'].astype(str).str.replace(r'\.0$', '', regex=True)

            return df_clean.dropna(subset=['Student ID', 'Obtained Marks'])

        try:
            xls_24   = pd.ExcelFile(file_24)
            df_24_bl = clean_sheet(pd.read_excel(file_24, sheet_name=0), 'AY24-25', 'Baseline')
            df_24_el = clean_sheet(
                pd.read_excel(file_24, sheet_name=1 if len(xls_24.sheet_names) > 1 else 0),
                'AY24-25', 'Endline',
            )
            xls_25   = pd.ExcelFile(file_25)
            df_25_bl = clean_sheet(pd.read_excel(file_25, sheet_name=0), 'AY25-26', 'Baseline')
            df_25_el = clean_sheet(
                pd.read_excel(file_25, sheet_name=1 if len(xls_25.sheet_names) > 1 else 0),
                'AY25-26', 'Endline',
            )
            return pd.concat([df_24_bl, df_24_el, df_25_bl, df_25_el], ignore_index=True)
        except Exception as e:
            st.error(f"Error loading multi-year data: {e}")
            return None

    FILE_24 = "EL-BL-Data-AY-24-25.xlsx"
    FILE_25 = "BL-EL-AY-25-26-Final-AllSubjects.xlsx"

    # ── Sidebar nav + optional file uploaders ──────────────────────
    with st.sidebar:
        try:
            st.image("evidyaloka_logo.png", width=273)
        except Exception:
            st.warning("⚠️ Logo not found.")
        st.success(f"👤 **Logged in as:** {st.session_state['user_first_name']}")
        nav_col1, nav_col2 = st.columns(2)
        with nav_col1:
            if st.button("🏠 Home", use_container_width=True, key="nav_home_long"):
                st.session_state["current_page"] = "home"
                st.rerun()
        with nav_col2:
            if st.button("Sign Out", use_container_width=True, key="signout_long"):
                for k in ["logged_in_email", "user_first_name"]:
                    st.session_state[k] = None if k == "logged_in_email" else "User"
                st.session_state["current_page"] = "home"
                st.rerun()
        st.markdown("---")
        st.info(
            "💡 **Module Note:** This section analyses the overlap between AY 24-25 "
            "and AY 25-26 to track long-term strategic growth."
        )

    # FIX 2 — file uploader fallback for longitudinal module
    # If the expected files are not on disk, offer upload widgets.
    src_24, src_25 = FILE_24, FILE_25
    if not (os.path.exists(FILE_24) and os.path.exists(FILE_25)):
        st.warning(
            f"⚠️ One or both data files (`{FILE_24}`, `{FILE_25}`) were not found on disk. "
            "Upload them below to continue."
        )
        up_24 = st.file_uploader(f"Upload AY 24-25 data ({FILE_24})", type=["xlsx"], key="up_24")
        up_25 = st.file_uploader(f"Upload AY 25-26 data ({FILE_25})", type=["xlsx"], key="up_25")
        if up_24 is None or up_25 is None:
            st.info("Please upload both files to load the Longitudinal Analysis module.")
            st.stop()
        src_24, src_25 = up_24, up_25

    with st.spinner("Synthesizing Multi-Year Intelligence..."):
        df_long = load_multi_year_data(src_24, src_25)

    filtered_df_long = pd.DataFrame()

    if df_long is not None and not df_long.empty:
        with st.sidebar:
            filtered_df_long, long_sel = build_filter_sidebar(df_long, key_prefix="long")

        if filtered_df_long.empty:
            st.warning("⚠️ No data available for the selected filters. Please adjust your criteria.")
        else:
            mig_tab, sub_tab, gen_tab, subj_tab_long, geo_tab_long, centre_tab_long = st.tabs([
                "📊 Overall Health (Migration)",
                "📚 Subject Efficacy",
                "🚻 Gender Equity",
                "📚 Subject Wise",
                "🗺️ Geographical Wise",
                "🏫 Centre Deep Dive",
            ])

            # Retained-cohort IDs (used by first 3 tabs only)
            df_el24 = filtered_df_long[filtered_df_long['Timepoint'] == 'AY24-25 Endline']
            df_el25 = filtered_df_long[filtered_df_long['Timepoint'] == 'AY25-26 Endline']

            if df_el24.empty or df_el25.empty:
                st.warning(
                    "⚠️ One or both endline years returned no rows after filtering. "
                    "Check whether the Subject or Grade filter is excluding an entire year. "
                    "Subject names may differ between workbooks (e.g. 'Math' vs 'Mathematics')."
                )

            ids_24 = set(df_el24['Student ID'].dropna().unique())
            ids_25 = set(df_el25['Student ID'].dropna().unique())
            retained_students = ids_24.intersection(ids_25)

            # Diagnostic expander — remove once overlap issue is resolved
            with st.expander("🛠️ Debug: ID Matching", expanded=False):
                dcol1, dcol2 = st.columns(2)
                with dcol1:
                    st.markdown("**AY 24-25 Endline**")
                    st.write(f"Rows: `{len(df_el24)}`")
                    st.write(f"Unique Student IDs: `{len(ids_24)}`")
                    sample_24 = (
                        df_el24['Student ID'].dropna().drop_duplicates().head(5).reset_index(drop=True)
                    )
                    st.dataframe(sample_24.rename("Student ID"), use_container_width=True)
                with dcol2:
                    st.markdown("**AY 25-26 Endline**")
                    st.write(f"Rows: `{len(df_el25)}`")
                    st.write(f"Unique Student IDs: `{len(ids_25)}`")
                    sample_25 = (
                        df_el25['Student ID'].dropna().drop_duplicates().head(5).reset_index(drop=True)
                    )
                    st.dataframe(sample_25.rename("Student ID"), use_container_width=True)
                st.markdown("---")
                st.write(f"**Intersection size:** `{len(retained_students)}`")
                if len(retained_students) == 0 and ids_24 and ids_25:
                    one_24 = next(iter(ids_24))
                    one_25 = next(iter(ids_25))
                    st.markdown("**Character-level inspection of first sample ID from each year:**")
                    st.write(f"AY24 → repr: `{repr(one_24)}` | len: `{len(one_24)}`")
                    st.write(f"AY25 → repr: `{repr(one_25)}` | len: `{len(one_25)}`")
                    st.caption(
                        "Look for: extra spaces, leading zeros (0001 vs 1), "
                        "Unicode non-breaking spaces (\\xa0), or BOM characters (\\ufeff). "
                        "repr() exposes all of these."
                    )

            no_overlap = len(retained_students) == 0
            df_ret_24  = df_el24[df_el24['Student ID'].isin(retained_students)] if not no_overlap else pd.DataFrame()
            df_ret_25  = df_el25[df_el25['Student ID'].isin(retained_students)] if not no_overlap else pd.DataFrame()

            # ── TAB 1: MIGRATION ───────────────────────────────────────
            with mig_tab:
                st.markdown("### 🧱 Structural Tier Migration (Retained Cohort)")
                if no_overlap:
                    st.warning(
                        "⚠️ No overlapping Student IDs were found between AY 24-25 and AY 25-26 "
                        "for the current filters. This tab requires matched students across both years. "
                        "Use the **Subject Wise**, **Geographical Wise**, or **Centre Deep Dive** tabs "
                        "for year-on-year comparison without requiring matched IDs."
                    )
                else:
                    col_m1, col_m2 = st.columns([1.5, 1])
                    with col_m1:
                        if 'Category' in df_ret_24.columns:
                            cat_24 = df_ret_24['Category'].value_counts(normalize=True).reset_index()
                            cat_24.columns = ['Category', 'Percentage']
                            cat_24['Year'] = 'AY 24-25 (Endline)'
                            cat_25 = df_ret_25['Category'].value_counts(normalize=True).reset_index()
                            cat_25.columns = ['Category', 'Percentage']
                            cat_25['Year'] = 'AY 25-26 (Endline)'
                            cat_df = pd.concat([cat_24, cat_25])
                            cat_df['Percentage'] *= 100
                            fig_cat = px.bar(
                                cat_df, x="Year", y="Percentage", color="Category",
                                color_discrete_map=RISE_COLORS_LONG,
                                text=cat_df['Percentage'].apply(lambda x: f'{x:.1f}%'),
                                category_orders={
                                    "Category": ["Reviving", "Initiating", "Shaping", "Evolving"],
                                    "Year": ["AY 24-25 (Endline)", "AY 25-26 (Endline)"],
                                },
                            )
                            fig_cat.update_layout(
                                barmode='stack', xaxis_title="", yaxis_title="% of Cohort",
                                plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30),
                                legend=dict(orientation="h", yanchor="top", y=-0.15, xanchor="center", x=0.5, title=""),
                            )
                            st.plotly_chart(fig_cat, width='stretch')
                    with col_m2:
                        st.info(
                            "**What this means:**\nThis chart strips away the noise and looks only at "
                            "students who stayed with us for two full years. A successful program will show "
                            "the red/orange sections shrinking as students migrate upward into green sections."
                        )
                        try:
                            rev_24 = cat_24[cat_24['Category'] == 'Reviving']['Percentage'].values
                            rev_25 = cat_25[cat_25['Category'] == 'Reviving']['Percentage'].values
                            rev_24_val = rev_24[0] * 100 if len(rev_24) else 0
                            rev_25_val = rev_25[0] * 100 if len(rev_25) else 0
                            rev_diff   = rev_25_val - rev_24_val
                            st.success(
                                f"**🔍 Key Insight:**\nThe proportion of critically struggling students "
                                f"('Reviving') changed by **{rev_diff:+.1f}%** Year-over-Year."
                            )
                            if rev_diff < 0:
                                st.markdown(
                                    "**💡 Strategic Suggestion:**\nExcellent progress. The base is shrinking. "
                                    "Keep investing in current foundational remedial strategies."
                                )
                            else:
                                st.markdown(
                                    "**💡 Strategic Suggestion:**\nThe struggling cohort is stagnating or growing. "
                                    "Consider implementing targeted small-group tutoring specifically for 'Reviving' students."
                                )
                        except Exception:
                            st.write("Insufficient category data for insights.")

            # ── TAB 2: SUBJECT EFFICACY ────────────────────────────────
            with sub_tab:
                st.markdown("### 📈 YoY Subject Trajectory (Slopegraph)")
                if no_overlap:
                    st.warning(
                        "⚠️ No overlapping Student IDs were found between AY 24-25 and AY 25-26. "
                        "This slopegraph requires matched students. Use the **Subject Wise** tab for a "
                        "full subject comparison across all timepoints without requiring matched IDs."
                    )
                else:
                    col_s1, col_s2 = st.columns([1.5, 1])
                    with col_s1:
                        subj_24 = df_ret_24.groupby('Subject')['Obtained Marks'].mean().reset_index()
                        subj_24['Year'] = 'AY 24-25'
                        subj_25 = df_ret_25.groupby('Subject')['Obtained Marks'].mean().reset_index()
                        subj_25['Year'] = 'AY 25-26'
                        slope_df = pd.concat([subj_24, subj_25])
                        fig_slope = px.line(
                            slope_df, x="Year", y="Obtained Marks", color="Subject",
                            markers=True, line_shape="linear", text="Obtained Marks",
                        )
                        fig_slope.update_traces(
                            textposition="top center", texttemplate='%{text:.1f}', marker=dict(size=10)
                        )
                        fig_slope.update_layout(
                            xaxis_title="", yaxis_title="Average Score", showlegend=True,
                            plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30),
                        )
                        fig_slope.update_xaxes(showgrid=False, linecolor='black')
                        fig_slope.update_yaxes(showgrid=True, gridcolor='lightgrey', zeroline=False)
                        st.plotly_chart(fig_slope, width='stretch')
                    with col_s2:
                        st.info(
                            "**What this means:**\nThis 'Slopegraph' visualizes momentum. The steepness "
                            "and direction of the lines reveal which subjects are improving and which are "
                            "backsliding over a 12-month period."
                        )
                        try:
                            growth_df  = pd.merge(subj_24, subj_25, on='Subject', suffixes=('_24', '_25'))
                            growth_df['Delta'] = growth_df['Obtained Marks_25'] - growth_df['Obtained Marks_24']
                            best_sub  = growth_df.loc[growth_df['Delta'].idxmax()]
                            worst_sub = growth_df.loc[growth_df['Delta'].idxmin()]
                            st.success(
                                f"**🔍 Key Insight:**\n**{best_sub['Subject']}** is the strongest performer, "
                                f"growing by {best_sub['Delta']:+.2f} points. **{worst_sub['Subject']}** showed "
                                f"the weakest momentum ({worst_sub['Delta']:+.2f} points)."
                            )
                            st.markdown(
                                f"**💡 Strategic Suggestion:**\nConduct a 'Bright Spot Analysis' on the "
                                f"{best_sub['Subject']} curriculum and cross-pollinate those techniques into "
                                f"{worst_sub['Subject']} planning for the upcoming semester."
                            )
                        except Exception:
                            st.write("Insufficient subject overlap for comparative insights.")

            # ── TAB 3: GENDER EQUITY ──────────────────────────────────
            with gen_tab:
                st.markdown("### 🚻 Gender Equity Tracking")
                if no_overlap:
                    st.warning(
                        "⚠️ No overlapping Student IDs were found between AY 24-25 and AY 25-26. "
                        "Gender equity tracking requires matched students across both years."
                    )
                elif 'Gender' in df_ret_24.columns and 'Gender' in df_ret_25.columns:
                    col_g1, col_g2 = st.columns([1.5, 1])
                    with col_g1:
                        g24 = (
                            df_ret_24[df_ret_24['Gender'].isin(['Boy', 'Girl'])]
                            .groupby('Gender')['Obtained Marks'].mean().reset_index()
                        )
                        g24['Year'] = 'AY 24-25'
                        g25 = (
                            df_ret_25[df_ret_25['Gender'].isin(['Boy', 'Girl'])]
                            .groupby('Gender')['Obtained Marks'].mean().reset_index()
                        )
                        g25['Year'] = 'AY 25-26'
                        gen_df = pd.concat([g24, g25])
                        fig_gen = px.bar(
                            gen_df, x="Year", y="Obtained Marks", color="Gender", barmode='group',
                            color_discrete_map={"Boy": "#636EFA", "Girl": "#EF553B"},
                            text=gen_df['Obtained Marks'].apply(lambda x: f'{x:.2f}'),
                        )
                        fig_gen.update_layout(
                            xaxis_title="", yaxis_title="Average Score",
                            plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30),
                            legend=dict(orientation="h", yanchor="top", y=-0.15, xanchor="center", x=0.5, title=""),
                        )
                        st.plotly_chart(fig_gen, width='stretch')
                    with col_g2:
                        st.info(
                            "**What this means:**\nThis metric tracks if our educational impact is equitable. "
                            "It ensures that as the overall average rises, one gender isn't being left behind."
                        )
                        try:
                            gap_24 = abs(
                                g24[g24['Gender'] == 'Girl']['Obtained Marks'].values[0]
                                - g24[g24['Gender'] == 'Boy']['Obtained Marks'].values[0]
                            )
                            gap_25 = abs(
                                g25[g25['Gender'] == 'Girl']['Obtained Marks'].values[0]
                                - g25[g25['Gender'] == 'Boy']['Obtained Marks'].values[0]
                            )
                            st.success(
                                f"**🔍 Key Insight:**\nThe performance gap between boys and girls was "
                                f"**{gap_24:.2f} points** last year, and is now **{gap_25:.2f} points** this year."
                            )
                            if gap_25 < gap_24:
                                st.markdown(
                                    "**💡 Strategic Suggestion:**\nThe equity gap is closing. Current inclusive "
                                    "classroom engagement strategies are working. Maintain the standard."
                                )
                            else:
                                st.markdown(
                                    "**💡 Strategic Suggestion:**\nThe equity gap is widening. Consider implementing "
                                    "targeted mentorship or gender-specific encouragement initiatives."
                                )
                        except Exception:
                            st.write("Insufficient gender data to calculate gaps.")
                else:
                    st.warning("Gender data is not consistently available across both academic years.")

            # ── TAB 4: SUBJECT WISE — Endline YoY only ────────────────
            with subj_tab_long:
                st.markdown("### 📚 Subject-Wise YoY Comparison (Endline Only)")
                df_endlines = filtered_df_long[filtered_df_long['Period'] == 'Endline'].copy()

                if 'Subject' in df_endlines.columns:
                    all_subjects = sorted(df_endlines['Subject'].dropna().unique())
                    selected_subj_long = st.selectbox(
                        "Select Subject", ["All"] + list(all_subjects), key="subj_long_select"
                    )
                    df_subj_view = df_endlines.copy()
                    if selected_subj_long != "All":
                        df_subj_view = df_subj_view[df_subj_view['Subject'].astype(str) == selected_subj_long]

                    df_subj_view['Academic Year'] = pd.Categorical(
                        df_subj_view['Academic Year'], categories=AY_ORDER, ordered=True
                    )

                    col_sw1, col_sw2 = st.columns(2)
                    with col_sw1:
                        st.markdown("#### Average Endline Score (YoY)")
                        subj_avg = (
                            df_subj_view
                            .groupby(['Academic Year', 'Subject'], observed=True)['Obtained Marks']
                            .mean().reset_index()
                        )
                        subj_avg = subj_avg.sort_values('Academic Year')
                        fig_sw_avg = px.line(
                            subj_avg, x="Academic Year", y="Obtained Marks", color="Subject",
                            markers=True, line_shape="linear", text="Obtained Marks",
                        )
                        fig_sw_avg.update_traces(
                            textposition="top center", texttemplate='%{text:.1f}', marker=dict(size=10)
                        )
                        fig_sw_avg.update_layout(
                            xaxis_title="", yaxis_title="Average Endline Score",
                            plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30),
                        )
                        fig_sw_avg.update_xaxes(showgrid=False, linecolor='black')
                        fig_sw_avg.update_yaxes(showgrid=True, gridcolor='lightgrey', zeroline=False)
                        st.plotly_chart(fig_sw_avg, width='stretch')

                    with col_sw2:
                        st.markdown("#### R.I.S.E Category Shift by Subject")
                        if 'Category' in df_subj_view.columns:
                            ay_options = [ay for ay in AY_ORDER if ay in df_subj_view['Academic Year'].values]
                            selected_ay = st.selectbox(
                                "Select Academic Year", ay_options,
                                index=len(ay_options) - 1, key="ay_subj_long",
                            )
                            df_subj_ay = df_subj_view[df_subj_view['Academic Year'] == selected_ay]
                            subj_cat = (
                                df_subj_ay
                                .groupby(['Subject', 'Category'], observed=True)
                                .size().reset_index(name='Count')
                            )
                            subj_cat['Percentage'] = (
                                subj_cat.groupby('Subject')['Count']
                                .transform(lambda x: x / x.sum() * 100)
                            )
                            fig_sw_rise = px.bar(
                                subj_cat, x="Subject", y="Percentage", color="Category",
                                color_discrete_map=RISE_COLORS_LONG,
                                text=subj_cat['Percentage'].apply(lambda x: f'{x:.1f}%' if x > 5 else ''),
                                category_orders={"Category": ["Reviving", "Initiating", "Shaping", "Evolving"]},
                            )
                            fig_sw_rise.update_layout(
                                barmode='stack', yaxis_title="% of Students", margin=dict(l=0, r=0, t=30),
                                legend=dict(orientation="h", yanchor="top", y=-0.2, xanchor="center", x=0.5, title=""),
                            )
                            st.plotly_chart(fig_sw_rise, width='stretch')

                    st.markdown("---")
                    st.markdown("#### Subject-Wise YoY Score Change (AY24-25 Endline → AY25-26 Endline)")
                    el_24_subj = (
                        df_endlines[df_endlines['Academic Year'] == 'AY24-25']
                        .groupby('Subject')['Obtained Marks'].mean()
                    )
                    el_25_subj = (
                        df_endlines[df_endlines['Academic Year'] == 'AY25-26']
                        .groupby('Subject')['Obtained Marks'].mean()
                    )
                    yoy_subj = (el_25_subj - el_24_subj).dropna().reset_index()
                    yoy_subj.columns = ['Subject', 'YoY Change']
                    yoy_subj['Color'] = yoy_subj['YoY Change'].apply(
                        lambda x: 'Improved' if x >= 0 else 'Declined'
                    )
                    if not yoy_subj.empty:
                        fig_yoy_subj = px.bar(
                            yoy_subj, x="Subject", y="YoY Change", color="Color",
                            color_discrete_map={"Improved": "#00964d", "Declined": "#ed1c2d"},
                            text=yoy_subj['YoY Change'].apply(lambda x: f'{x:+.2f}'),
                        )
                        fig_yoy_subj.add_hline(y=0, line_dash="dash", line_color="black")
                        fig_yoy_subj.update_layout(
                            showlegend=False, margin=dict(l=0, r=0, t=30),
                            plot_bgcolor="rgba(0,0,0,0)", xaxis_title="",
                            yaxis_title="YoY Endline Score Change",
                        )
                        st.plotly_chart(fig_yoy_subj, width='stretch')
                    else:
                        st.info(
                            "Both AY24-25 and AY25-26 Endline data are needed for this comparison. "
                            "Check that current filters do not exclude one academic year entirely."
                        )
                else:
                    st.warning("Subject data not available.")

            # ── TAB 5: GEOGRAPHICAL WISE — Endline YoY only ───────────
            with geo_tab_long:
                st.markdown("### 🗺️ Geographical YoY Comparison (Endline Only)")
                df_geo = filtered_df_long[filtered_df_long['Period'] == 'Endline'].copy()

                if 'State' in df_geo.columns:
                    df_geo['Academic Year'] = pd.Categorical(
                        df_geo['Academic Year'], categories=AY_ORDER, ordered=True
                    )

                    col_g1, col_g2 = st.columns(2)
                    with col_g1:
                        st.markdown("#### State-wise Average Endline Score (YoY)")
                        state_avg = (
                            df_geo.groupby(['State', 'Academic Year'], observed=True)['Obtained Marks']
                            .mean().reset_index()
                        )
                        fig_state_avg = px.line(
                            state_avg, x="Academic Year", y="Obtained Marks", color="State",
                            markers=True, line_shape="linear", text="Obtained Marks",
                        )
                        fig_state_avg.update_traces(
                            textposition="top center", texttemplate='%{text:.1f}', marker=dict(size=9)
                        )
                        fig_state_avg.update_layout(
                            xaxis_title="", yaxis_title="Average Endline Score",
                            plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30),
                        )
                        fig_state_avg.update_xaxes(showgrid=False, linecolor='black')
                        fig_state_avg.update_yaxes(showgrid=True, gridcolor='lightgrey', zeroline=False)
                        st.plotly_chart(fig_state_avg, width='stretch')

                    with col_g2:
                        st.markdown("#### State-wise R.I.S.E Shift")
                        if 'Category' in df_geo.columns:
                            ay_opts_geo = [ay for ay in AY_ORDER if ay in df_geo['Academic Year'].values]
                            selected_ay_geo = st.selectbox(
                                "Select Academic Year", ay_opts_geo,
                                index=len(ay_opts_geo) - 1, key="ay_geo_long",
                            )
                            state_cat_filtered = df_geo[df_geo['Academic Year'] == selected_ay_geo]
                            state_cat_grp = (
                                state_cat_filtered.groupby(['State', 'Category'])
                                .size().reset_index(name='Count')
                            )
                            state_cat_grp['Percentage'] = (
                                state_cat_grp.groupby('State')['Count']
                                .transform(lambda x: x / x.sum() * 100)
                            )
                            fig_state_rise = px.bar(
                                state_cat_grp, x="State", y="Percentage", color="Category",
                                color_discrete_map=RISE_COLORS_LONG,
                                text=state_cat_grp['Percentage'].apply(lambda x: f'{x:.1f}%' if x > 5 else ''),
                                category_orders={"Category": ["Reviving", "Initiating", "Shaping", "Evolving"]},
                            )
                            fig_state_rise.update_layout(
                                barmode='stack', yaxis_title="% of Students", margin=dict(l=0, r=0, t=30),
                                legend=dict(orientation="h", yanchor="top", y=-0.2, xanchor="center", x=0.5, title=""),
                            )
                            st.plotly_chart(fig_state_rise, width='stretch')

                    st.markdown("---")
                    st.markdown("#### State-wise YoY Score Change (AY24-25 Endline → AY25-26 Endline)")
                    el24_state = (
                        df_geo[df_geo['Academic Year'] == 'AY24-25']
                        .groupby('State')['Obtained Marks'].mean()
                    )
                    el25_state = (
                        df_geo[df_geo['Academic Year'] == 'AY25-26']
                        .groupby('State')['Obtained Marks'].mean()
                    )
                    if not el24_state.empty and not el25_state.empty:
                        yoy_state = (el25_state - el24_state).dropna().reset_index()
                        yoy_state.columns = ['State', 'YoY Change']
                        yoy_state['Color'] = yoy_state['YoY Change'].apply(
                            lambda x: 'Improved' if x >= 0 else 'Declined'
                        )
                        fig_yoy_state = px.bar(
                            yoy_state, x="State", y="YoY Change", color="Color",
                            color_discrete_map={"Improved": "#00964d", "Declined": "#ed1c2d"},
                            text=yoy_state['YoY Change'].apply(lambda x: f'{x:+.2f}'),
                        )
                        fig_yoy_state.add_hline(y=0, line_dash="dash", line_color="black")
                        fig_yoy_state.update_layout(
                            showlegend=False, margin=dict(l=0, r=0, t=30),
                            plot_bgcolor="rgba(0,0,0,0)", xaxis_title="",
                            yaxis_title="YoY Endline Score Change",
                        )
                        st.plotly_chart(fig_yoy_state, width='stretch')
                    else:
                        st.info("Both AY24-25 and AY25-26 Endline data are needed for YoY state comparison.")
                else:
                    st.warning("State data not available.")

            # ── TAB 6: CENTRE DEEP DIVE — Endline YoY only ────────────
            with centre_tab_long:
                st.markdown("### 🏫 Centre Deep Dive — Individual Centre YoY Comparison (Endline Only)")
                df_cdd = filtered_df_long[filtered_df_long['Period'] == 'Endline'].copy()

                if 'Centre Name' in df_cdd.columns:
                    df_cdd['Academic Year'] = pd.Categorical(
                        df_cdd['Academic Year'], categories=AY_ORDER, ordered=True
                    )
                    all_centres_long = sorted(df_cdd['Centre Name'].dropna().unique())
                    selected_centre_deep = st.selectbox(
                        "Select Centre for Deep Dive", all_centres_long, key="centre_deep_long"
                    )
                    df_centre_deep = df_cdd[df_cdd['Centre Name'].astype(str) == selected_centre_deep]

                    if not df_centre_deep.empty:
                        col_c1, col_c2 = st.columns(2)
                        with col_c1:
                            st.markdown("#### Average Endline Score Trajectory (YoY)")
                            centre_avg = (
                                df_centre_deep
                                .groupby(['Academic Year', 'Subject'], observed=True)['Obtained Marks']
                                .mean().reset_index()
                            )
                            fig_c_avg = px.line(
                                centre_avg, x="Academic Year", y="Obtained Marks", color="Subject",
                                markers=True, line_shape="linear", text="Obtained Marks",
                            )
                            fig_c_avg.update_traces(
                                textposition="top center", texttemplate='%{text:.1f}', marker=dict(size=10)
                            )
                            fig_c_avg.update_layout(
                                xaxis_title="", yaxis_title="Average Endline Score",
                                plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30),
                            )
                            fig_c_avg.update_xaxes(showgrid=False, linecolor='black')
                            fig_c_avg.update_yaxes(showgrid=True, gridcolor='lightgrey', zeroline=False)
                            st.plotly_chart(fig_c_avg, width='stretch')

                        with col_c2:
                            st.markdown("#### R.I.S.E Distribution by Academic Year")
                            if 'Category' in df_centre_deep.columns:
                                centre_cat = (
                                    df_centre_deep
                                    .groupby(['Academic Year', 'Category'], observed=True)
                                    .size().reset_index(name='Count')
                                )
                                centre_cat['Percentage'] = (
                                    centre_cat.groupby('Academic Year')['Count']
                                    .transform(lambda x: x / x.sum() * 100)
                                )
                                fig_c_rise = px.bar(
                                    centre_cat, x="Academic Year", y="Percentage", color="Category",
                                    color_discrete_map=RISE_COLORS_LONG,
                                    text=centre_cat['Percentage'].apply(lambda x: f'{x:.1f}%' if x > 5 else ''),
                                    category_orders={
                                        "Category": ["Reviving", "Initiating", "Shaping", "Evolving"],
                                        "Academic Year": AY_ORDER,
                                    },
                                )
                                fig_c_rise.update_layout(
                                    barmode='stack', yaxis_title="% of Students", margin=dict(l=0, r=0, t=30),
                                    legend=dict(orientation="h", yanchor="top", y=-0.2, xanchor="center", x=0.5, title=""),
                                )
                                st.plotly_chart(fig_c_rise, width='stretch')

                        st.markdown("---")
                        st.markdown("#### Centre vs All-Centre Benchmark (Endlines Only, YoY)")
                        el24_centre = df_centre_deep[df_centre_deep['Academic Year'] == 'AY24-25']['Obtained Marks'].mean()
                        el25_centre = df_centre_deep[df_centre_deep['Academic Year'] == 'AY25-26']['Obtained Marks'].mean()
                        el24_all    = df_cdd[df_cdd['Academic Year'] == 'AY24-25']['Obtained Marks'].mean()
                        el25_all    = df_cdd[df_cdd['Academic Year'] == 'AY25-26']['Obtained Marks'].mean()
                        bench_data  = pd.DataFrame({
                            'Group':        [selected_centre_deep, 'All Centres', selected_centre_deep, 'All Centres'],
                            'Academic Year': ['AY24-25', 'AY24-25', 'AY25-26', 'AY25-26'],
                            'Avg Score':    [el24_centre, el24_all, el25_centre, el25_all],
                        }).dropna(subset=['Avg Score'])
                        if not bench_data.empty:
                            bench_data['Academic Year'] = pd.Categorical(
                                bench_data['Academic Year'], categories=AY_ORDER, ordered=True
                            )
                            fig_bench = px.bar(
                                bench_data, x="Academic Year", y="Avg Score", color="Group",
                                barmode='group',
                                text=bench_data['Avg Score'].apply(lambda x: f'{x:.2f}'),
                            )
                            fig_bench.update_layout(
                                xaxis_title="", yaxis_title="Average Endline Score",
                                margin=dict(l=0, r=0, t=30),
                                legend=dict(orientation="h", yanchor="top", y=-0.15, xanchor="center", x=0.5, title=""),
                            )
                            st.plotly_chart(fig_bench, width='stretch')

                        st.markdown("---")
                        st.markdown("#### Gender Breakdown at This Centre (Endline YoY)")
                        if 'Gender' in df_centre_deep.columns:
                            gen_centre = df_centre_deep[
                                ~df_centre_deep['Gender'].astype(str).str.lower().isin(['nan', 'none', 'null', ''])
                            ]
                            if not gen_centre.empty:
                                gen_avg_c = (
                                    gen_centre
                                    .groupby(['Academic Year', 'Gender'], observed=True)['Obtained Marks']
                                    .mean().reset_index()
                                )
                                gen_avg_c['Academic Year'] = pd.Categorical(
                                    gen_avg_c['Academic Year'], categories=AY_ORDER, ordered=True
                                )
                                fig_gen_c = px.line(
                                    gen_avg_c, x="Academic Year", y="Obtained Marks", color="Gender",
                                    markers=True, line_shape="linear",
                                    color_discrete_map={"Boy": "#636EFA", "Girl": "#EF553B"},
                                )
                                fig_gen_c.update_traces(marker=dict(size=10))
                                fig_gen_c.update_layout(
                                    xaxis_title="", yaxis_title="Average Endline Score",
                                    plot_bgcolor="rgba(0,0,0,0)", margin=dict(l=0, r=0, t=30),
                                )
                                st.plotly_chart(fig_gen_c, width='stretch')
                            else:
                                st.info("No valid gender data for this centre.")
                    else:
                        st.warning("No Endline data found for this centre with the current filters.")
                else:
                    st.warning("Centre Name data not available.")
    else:
        st.error(
            f"⚠️ **Longitudinal Data Error!** Could not load data from the provided files. "
            f"Please check that `{FILE_24}` and `{FILE_25}` are valid Excel workbooks."
        )

    st.stop()


# ==========================================
# MAIN DASHBOARD
# ==========================================
st.title("📈 Impact Analytics Dashboard")
st.markdown(
    "<p style='color:gray;font-size:1.1em;'>Comprehensive Baseline vs. Endline Performance Assessment</p>",
    unsafe_allow_html=True,
)

with st.sidebar:
    try:
        st.image("evidyaloka_logo.png", width=273)
    except Exception:
        st.warning("⚠️ Logo not found.")
    st.success(f"👤 **Logged in as:** {st.session_state['user_first_name']}")
    nav_col1, nav_col2 = st.columns(2)
    with nav_col1:
        if st.button("🏠 Home", use_container_width=True, key="nav_home_main"):
            st.session_state["current_page"] = "home"
            st.rerun()
    with nav_col2:
        if st.button("Sign Out", use_container_width=True, key="signout_main"):
            for k in ["logged_in_email", "user_first_name"]:
                st.session_state[k] = None if k == "logged_in_email" else "User"
            st.session_state["current_page"] = "home"
            st.rerun()
    st.markdown("---")


# ==========================================
# DATA LOADING ENGINE
# ==========================================
@st.cache_data
def load_and_prep_data(file_source):
    common_cols = [
        'State', 'Centre Name', 'Donor', 'Subject', 'Grade',
        'Student ID', 'Gender', 'Total Marks', 'Obtained Marks', 'Category', 'Academic Year',
    ]
    dfs_to_concat = []
    try:
        xls         = pd.ExcelFile(file_source)
        sheet_names = xls.sheet_names
        base_sheet  = 'Baseline' if 'Baseline' in sheet_names else 0
        end_sheet   = 'Endline'  if 'Endline'  in sheet_names else (1 if len(sheet_names) > 1 else 0)

        df_base = pd.read_excel(file_source, sheet_name=base_sheet)
        if 'Rubrics' in df_base.columns:
            df_base.rename(columns={'Rubrics': 'Category'}, inplace=True)
        df_base['Academic Year'] = 'Baseline'
        dfs_to_concat.append(df_base[[c for c in common_cols if c in df_base.columns]])

        df_end = pd.read_excel(file_source, sheet_name=end_sheet)
        if 'Rubrics' in df_end.columns:
            df_end.rename(columns={'Rubrics': 'Category'}, inplace=True)
        df_end['Academic Year'] = 'Endline'
        dfs_to_concat.append(df_end[[c for c in common_cols if c in df_end.columns]])
    except Exception as e:
        st.error(f"Error reading the Excel file: {e}")
        return pd.DataFrame()

    if not dfs_to_concat:
        return pd.DataFrame()

    df_combined = pd.concat(dfs_to_concat, ignore_index=True)
    for col in ['State', 'Centre Name', 'Donor', 'Subject', 'Student ID', 'Gender', 'Category']:
        if col in df_combined.columns:
            df_combined[col] = df_combined[col].astype(str).str.strip()
    if 'Gender' in df_combined.columns:
        df_combined['Gender'] = df_combined['Gender'].astype(str).str.strip().str.title()
    if 'Grade' in df_combined.columns:
        df_combined['Grade'] = df_combined['Grade'].astype(str).str.replace(r'\.0$', '', regex=True)
    if 'Category' in df_combined.columns:
        rise_order = ["Reviving", "Initiating", "Shaping", "Evolving"]
        df_combined['Category'] = pd.Categorical(df_combined['Category'], categories=rise_order, ordered=True)
    df_combined['Obtained Marks'] = pd.to_numeric(df_combined['Obtained Marks'], errors='coerce')
    return df_combined


DATA_FILE = "BL-EL-AY-25-26-Final-AllSubjects.xlsx"

# FIX 2 — file uploader fallback for main dashboard
data_source = None
if os.path.exists(DATA_FILE):
    data_source = DATA_FILE
else:
    st.warning(
        f"⚠️ `{DATA_FILE}` was not found on disk. Upload it below to load the dashboard."
    )
    uploaded_main = st.file_uploader(
        f"Upload AY 25-26 data ({DATA_FILE})", type=["xlsx"], key="up_main"
    )
    if uploaded_main is not None:
        data_source = uploaded_main

if data_source is not None:
    with st.spinner('Loading and crunching numbers...'):
        df = load_and_prep_data(data_source)

    if not df.empty:
        with st.sidebar:
            filtered_df, main_sel = build_filter_sidebar(df, key_prefix="main")
        selected_donors = main_sel["donor"]

        if filtered_df.empty:
            st.warning("⚠️ No data available for the selected filters. Please adjust your criteria.")
        else:
            tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs([
                "📊 Executive Summary",
                "📚 Subject Deep-Dive",
                "🗺️ Geographic View",
                "🧑‍🎓 Student-Level Impact",
                "🚻 Gender Analysis",
                "📉 RTM Analysis",
            ])

            base_df = filtered_df[filtered_df['Academic Year'] == 'Baseline']
            end_df  = filtered_df[filtered_df['Academic Year'] == 'Endline']

            # ------------------------------------------
            # TAB 1: EXECUTIVE SUMMARY
            # ------------------------------------------
            with tab1:
                st.markdown("### 🚀 High-Level Metrics")
                kpi1, kpi2, kpi3, kpi4, kpi5 = st.columns(5)
                if not base_df.empty and not end_df.empty and 'Student ID' in df.columns:
                    matched_students = len(
                        pd.merge(
                            base_df[['Student ID', 'Subject']].dropna(),
                            end_df[['Student ID', 'Subject']].dropna(),
                            on=['Student ID', 'Subject'],
                        )
                    )
                else:
                    matched_students = 0

                avg_base = base_df['Obtained Marks'].mean() if not base_df.empty else None
                avg_end  = end_df['Obtained Marks'].mean()  if not end_df.empty  else None
                sd_base  = base_df['Obtained Marks'].std()  if not base_df.empty and len(base_df) > 1 else None
                sd_end   = end_df['Obtained Marks'].std()   if not end_df.empty  and len(end_df) > 1  else None

                kpi1.metric("Matched Students", f"{matched_students:,}")
                if avg_base is not None and avg_end is not None:
                    kpi2.metric("Baseline Mean Score", f"{avg_base:.2f}")
                    kpi3.metric("Endline Mean Score",  f"{avg_end:.2f}",  delta=f"{avg_end - avg_base:.2f}")
                    if sd_base is not None and sd_end is not None:
                        kpi4.metric("Endline Score SD", f"{sd_end:.2f}", delta=f"{sd_end - sd_base:.2f}", delta_color="inverse")
                    else:
                        kpi4.metric("Endline Score SD", "N/A")
                    base_evolve = len(base_df[base_df['Category'] == 'Evolving']) / len(base_df) * 100 if len(base_df) > 0 else 0
                    end_evolve  = len(end_df[end_df['Category']   == 'Evolving']) / len(end_df)  * 100 if len(end_df)  > 0 else 0
                    kpi5.metric("Students in 'Evolving'", f"{end_evolve:.1f}%", delta=f"{end_evolve - base_evolve:.1f}%")
                elif avg_base is not None:
                    kpi2.metric("Baseline Mean Score", f"{avg_base:.2f}")
                    kpi3.metric("Endline Mean Score", "N/A")
                    kpi4.metric("Endline Score SD", "N/A")
                    kpi5.metric("Data Status", "Awaiting Endline")
                else:
                    kpi2.metric("Baseline Mean Score", "N/A")
                    kpi3.metric("Endline Mean Score", f"{avg_end:.2f}" if avg_end else "N/A")
                    kpi4.metric("Endline Score SD", f"{sd_end:.2f}" if sd_end else "N/A")
                    kpi5.metric("Data Status", "Endline Only")

                st.info(
                    "**💡 Understanding Standard Deviation (SD):** SD measures the spread of student scores. "
                    "A **decrease** (green delta) in SD means scores are becoming more clustered, indicating "
                    "the learning gap between high and low performers is closing. An **increase** (red delta) "
                    "means the gap is widening."
                )
                st.markdown("---")
                col_a, col_b = st.columns(2)
                with col_a:
                    st.markdown("#### 📈 Score Distribution (Box Plot)")
                    st.caption("Visualizes the spread of scores, median, and outliers.")
                    fig_box = px.box(
                        filtered_df, x="Academic Year", y="Obtained Marks",
                        color="Academic Year", color_discrete_map=COLOR_MAP, points="all",
                    )
                    fig_box.update_layout(showlegend=False, margin=dict(l=0, r=0, t=30))
                    st.plotly_chart(fig_box, width='stretch')
                with col_b:
                    st.markdown("#### 🧬 R.I.S.E Category Shift")
                    st.caption("Proportional breakdown of performance categories.")
                    cat_counts = filtered_df.groupby(['Academic Year', 'Category']).size().reset_index(name='Count')
                    cat_counts['Percentage'] = cat_counts.groupby('Academic Year')['Count'].transform(
                        lambda x: x / x.sum() * 100
                    )
                    cat_counts = cat_counts.sort_values(['Academic Year', 'Category'])
                    fig_rise = px.bar(
                        cat_counts, x="Category", y="Percentage", color="Academic Year",
                        text=cat_counts['Percentage'].apply(lambda x: f'{x:.1f}%' if not pd.isna(x) else ''),
                        color_discrete_map=COLOR_MAP,
                        category_orders={"Category": ["Reviving", "Initiating", "Shaping", "Evolving"]},
                    )
                    fig_rise.update_layout(
                        barmode='group', margin=dict(l=0, r=0, t=30),
                        legend=dict(orientation="h", yanchor="top", y=-0.15, xanchor="center", x=0.5, title=""),
                    )
                    st.plotly_chart(fig_rise, width='stretch')

            # ------------------------------------------
            # TAB 2: SUBJECT DEEP-DIVE
            # ------------------------------------------
            with tab2:
                st.markdown("### 📚 Subject & Grade Performance (R.I.S.E. Distribution)")

                def get_stacked_data(df_subset):
                    if df_subset.empty or 'Grade' not in df_subset.columns:
                        return pd.DataFrame()
                    grouped = df_subset.groupby(['Grade', 'Category']).size().reset_index(name='Count')
                    grouped['Percentage'] = grouped.groupby('Grade')['Count'].transform(
                        lambda x: x / x.sum() * 100
                    )
                    return grouped

                base_stacked = get_stacked_data(base_df)
                end_stacked  = get_stacked_data(end_df)
                sub_col1, sub_col2 = st.columns(2)
                with sub_col1:
                    st.markdown("#### Baseline R.I.S.E by Grade")
                    if not base_stacked.empty:
                        fig_base_grade = px.bar(
                            base_stacked, x="Grade", y="Percentage", color="Category",
                            color_discrete_map=RISE_COLORS,
                            text=base_stacked['Percentage'].apply(lambda x: f'{x:.1f}%' if x > 5 else ''),
                            category_orders={"Category": ["Reviving", "Initiating", "Shaping", "Evolving"]},
                        )
                        fig_base_grade.update_layout(
                            barmode='stack', yaxis_title="% of Students", margin=dict(l=0, r=0, t=30),
                            legend=dict(orientation="h", yanchor="top", y=-0.15, xanchor="center", x=0.5, title=""),
                        )
                        st.plotly_chart(fig_base_grade, width='stretch')
                    else:
                        st.info("No Baseline data available.")
                with sub_col2:
                    st.markdown("#### Endline R.I.S.E by Grade")
                    if not end_stacked.empty:
                        fig_end_grade = px.bar(
                            end_stacked, x="Grade", y="Percentage", color="Category",
                            color_discrete_map=RISE_COLORS,
                            text=end_stacked['Percentage'].apply(lambda x: f'{x:.1f}%' if x > 5 else ''),
                            category_orders={"Category": ["Reviving", "Initiating", "Shaping", "Evolving"]},
                        )
                        fig_end_grade.update_layout(
                            barmode='stack', yaxis_title="% of Students", margin=dict(l=0, r=0, t=30),
                            legend=dict(orientation="h", yanchor="top", y=-0.15, xanchor="center", x=0.5, title=""),
                        )
                        st.plotly_chart(fig_end_grade, width='stretch')
                    else:
                        st.info("No Endline data available.")

                st.markdown("---")
                st.markdown("#### 🧠 Automated Insights")
                if not base_stacked.empty and not end_stacked.empty:
                    try:
                        base_piv = base_stacked.pivot(index='Grade', columns='Category', values='Percentage').fillna(0)
                        end_piv  = end_stacked.pivot(index='Grade',  columns='Category', values='Percentage').fillna(0)
                        for cat in ["Reviving", "Initiating", "Shaping", "Evolving"]:
                            if cat not in base_piv.columns: base_piv[cat] = 0
                            if cat not in end_piv.columns:  end_piv[cat]  = 0
                        common_grades = base_piv.index.intersection(end_piv.index)
                        if len(common_grades) > 0:
                            diff_piv = end_piv.loc[common_grades] - base_piv.loc[common_grades]
                            best_evo_grade = diff_piv['Evolving'].idxmax()
                            best_evo_val   = diff_piv['Evolving'].max()
                            best_rev_grade = diff_piv['Reviving'].idxmin()
                            best_rev_val   = diff_piv['Reviving'].min()
                            if best_evo_val > 0:
                                st.success(
                                    f"📈 **Top Excellence Growth:** Grade **{best_evo_grade}** saw the highest shift "
                                    f"into 'Evolving', increasing its top-tier share by **{best_evo_val:+.1f}** percentage points."
                                )
                            else:
                                st.warning("⚠️ **Excellence Alert:** No grade saw an increase in the 'Evolving' category percentage.")
                            if best_rev_val < 0:
                                st.success(
                                    f"📉 **Highest Risk Reduction:** Grade **{best_rev_grade}** had the most successful intervention "
                                    f"for struggling students, reducing its 'Reviving' population by **{abs(best_rev_val):.1f}** percentage points."
                                )
                            else:
                                st.warning("⚠️ **Risk Alert:** No grade successfully reduced their share of students in the 'Reviving' category.")
                        else:
                            st.info("Insufficient overlapping grades between Baseline and Endline to generate comparative insights.")
                    except Exception:
                        st.info("Not enough data variance to generate automated insights.")
                else:
                    st.info("Awaiting both Baseline and Endline data to generate comparative insights.")

            # ------------------------------------------
            # TAB 3: GEOGRAPHIC VIEW
            # ------------------------------------------
            with tab3:
                st.markdown("### 🗺️ Geographic & Centre Analysis")
                st.markdown("#### State-wise Performance Comparison (R.I.S.E %)")
                if not filtered_df.empty:
                    state_cat = filtered_df.groupby(['State', 'Academic Year', 'Category']).size().reset_index(name='Count')
                    state_cat['Percentage'] = state_cat.groupby(['State', 'Academic Year'])['Count'].transform(
                        lambda x: x / x.sum() * 100
                    )
                    state_cat['Period'] = state_cat['Academic Year'].map({'Baseline': 'B', 'Endline': 'E'})

                    def abbreviate_state(s):
                        words = str(s).split()
                        return "".join(w.upper() for w in words) if len(words) > 1 else str(s)[:3].upper()

                    state_cat['State Abbr'] = state_cat['State'].apply(abbreviate_state)
                    fig_state = px.bar(
                        state_cat, x="Period", y="Percentage", color="Category", facet_col="State Abbr",
                        hover_data={"State": True, "State Abbr": False, "Period": False, "Academic Year": True},
                        color_discrete_map=RISE_COLORS,
                        text=state_cat['Percentage'].apply(lambda x: f'{x:.1f}%' if not pd.isna(x) and x > 5 else ''),
                        category_orders={
                            "Category": ["Reviving", "Initiating", "Shaping", "Evolving"],
                            "Period": ["B", "E"],
                        },
                    )
                    fig_state.update_layout(
                        barmode='stack', yaxis_title="% of Students", margin=dict(l=0, r=0, t=40),
                        legend=dict(orientation="h", yanchor="top", y=-0.2, xanchor="center", x=0.5, title=""),
                    )
                    fig_state.update_xaxes(title_text='')
                    fig_state.for_each_annotation(lambda a: a.update(text=a.text.split("=")[-1]))
                    st.plotly_chart(fig_state, width='stretch')
                else:
                    st.info("No data available for State comparison.")

                st.markdown("---")
                st.markdown("#### Top 10 Centres (Sorted by % Evolving)")
                if not filtered_df.empty:
                    center_cat = filtered_df.groupby(['Centre Name', 'Category']).size().reset_index(name='Count')
                    center_cat['Percentage'] = center_cat.groupby('Centre Name')['Count'].transform(
                        lambda x: x / x.sum() * 100
                    )
                    center_piv = center_cat.pivot(index='Centre Name', columns='Category', values='Percentage').fillna(0)
                    for cat in ["Reviving", "Initiating", "Shaping", "Evolving"]:
                        if cat not in center_piv.columns:
                            center_piv[cat] = 0
                    center_piv_sorted = center_piv.sort_values(
                        by=['Evolving', 'Shaping', 'Initiating', 'Reviving'],
                        ascending=[False, False, False, False],
                    ).head(10).iloc[::-1]
                    top_centres_long = center_piv_sorted.reset_index().melt(
                        id_vars='Centre Name',
                        value_vars=["Reviving", "Initiating", "Shaping", "Evolving"],
                        var_name='Category', value_name='Percentage',
                    )
                    fig_top_centres = px.bar(
                        top_centres_long, x="Percentage", y="Centre Name", color="Category",
                        orientation='h', color_discrete_map=RISE_COLORS,
                        text=top_centres_long['Percentage'].apply(lambda x: f'{x:.1f}%' if x > 5 else ''),
                        category_orders={"Category": ["Reviving", "Initiating", "Shaping", "Evolving"]},
                    )
                    fig_top_centres.update_layout(
                        barmode='stack', xaxis_title="% of Students", yaxis_title="",
                        margin=dict(l=0, r=0, t=30),
                        legend=dict(orientation="h", yanchor="top", y=-0.15, xanchor="center", x=0.5, title=""),
                    )
                    st.plotly_chart(fig_top_centres, width='stretch')
                else:
                    st.info("No data available for Top Centres.")

            # ------------------------------------------
            # TAB 4: STUDENT-LEVEL IMPACT
            # ------------------------------------------
            with tab4:
                st.markdown("### 🧑‍🎓 Student-Level Impact (Matched Cohort)")
                st.markdown("Tracking individual student growth by matching their Baseline and Endline records.")
                if not base_df.empty and not end_df.empty and 'Student ID' in df.columns:
                    base_clean = base_df[['Student ID', 'Subject', 'Obtained Marks', 'Category']].dropna(subset=['Student ID'])
                    end_clean  = end_df[['Student ID',  'Subject', 'Obtained Marks', 'Category']].dropna(subset=['Student ID'])
                    base_clean = base_clean.drop_duplicates(subset=['Student ID', 'Subject'])
                    end_clean  = end_clean.drop_duplicates(subset=['Student ID', 'Subject'])
                    paired_df  = pd.merge(base_clean, end_clean, on=['Student ID', 'Subject'], suffixes=('_BL', '_EL'))

                    if not paired_df.empty:
                        paired_df['Score Delta'] = paired_df['Obtained Marks_EL'] - paired_df['Obtained Marks_BL']
                        mean_change  = paired_df['Score Delta'].mean()
                        total_paired = len(paired_df)
                        positive_pct = len(paired_df[paired_df['Score Delta'] > 0]) / total_paired * 100
                        neutral_pct  = len(paired_df[paired_df['Score Delta'] == 0]) / total_paired * 100
                        negative_pct = len(paired_df[paired_df['Score Delta'] < 0]) / total_paired * 100

                        st.markdown("---")
                        m1, m2, m3, m4, m5 = st.columns(5)
                        m1.metric("Matched Students",     f"{total_paired:,}")
                        m2.metric("Avg Score Change",     f"{mean_change:+.2f}")
                        m3.metric("Students (+ Score)",   f"{positive_pct:.1f}%")
                        m4.metric("Students (No Change)", f"{neutral_pct:.1f}%")
                        m5.metric("Students (- Score)",   f"{negative_pct:.1f}%")
                        st.markdown("---")
                        st.markdown("#### 🔄 Category Transition Matrix")
                        st.caption(
                            "Read rows left-to-right to see student mobility. "
                            "**Background colors represent transition status:** "
                            "<span style='color:#82E0AA;font-weight:bold;'>Green (Upward)</span>, "
                            "<span style='color:#A9A9A9;font-weight:bold;'>Grey (No Change)</span>, and "
                            "<span style='color:#FF7F7F;font-weight:bold;'>Red (Downward)</span>.",
                            unsafe_allow_html=True,
                        )
                        transition = pd.crosstab(
                            paired_df['Category_BL'], paired_df['Category_EL'], normalize='index'
                        ) * 100
                        cat_order = ["Reviving", "Initiating", "Shaping", "Evolving"]
                        transition = transition.reindex(index=cat_order, columns=cat_order, fill_value=0)
                        direction_matrix = pd.DataFrame(index=cat_order, columns=cat_order)
                        for i, bl in enumerate(cat_order):
                            for j, el in enumerate(cat_order):
                                direction_matrix.loc[bl, el] = 0 if i == j else (1 if j > i else -1)
                        direction_matrix = direction_matrix.astype(float)
                        fig_heat = px.imshow(
                            direction_matrix,
                            labels=dict(x="Endline Category", y="Baseline Category", color="Transition Type"),
                            x=transition.columns, y=transition.index,
                            color_continuous_scale=["#FF7F7F", "#F2F4F7", "#82E0AA"],
                        )
                        text_matrix = transition.map(lambda x: f"{x:.1f}%")
                        fig_heat.update_traces(
                            text=text_matrix, texttemplate="%{text}",
                            hovertemplate="Baseline: %{y}<br>Endline: %{x}<br>Students: %{text}<extra></extra>",
                        )
                        fig_heat.update_coloraxes(showscale=False)
                        fig_heat.update_layout(margin=dict(l=0, r=0, t=30, b=0), height=500)
                        _, c2, _ = st.columns(3)
                        with c2:
                            st.plotly_chart(fig_heat, width='stretch')
                    else:
                        st.warning("⚠️ Could not find matching 'Student ID' and 'Subject' between the Baseline and Endline datasets.")
                else:
                    st.info("⚠️ Both Baseline and Endline datasets with a valid 'Student ID' column are required for this analysis.")

            # ------------------------------------------
            # TAB 5: GENDER ANALYSIS
            # ------------------------------------------
            with tab5:
                st.markdown("### 🚻 Gender-Wise Performance")
                if 'Gender' in filtered_df.columns:
                    gdf = filtered_df[
                        ~filtered_df['Gender'].astype(str).str.lower().isin(['nan', 'none', 'null', ''])
                    ].copy()
                    if not gdf.empty:
                        st.markdown("#### 🏆 Endline Average Score Snapshot")
                        g_base = gdf[gdf['Academic Year'] == 'Baseline']
                        g_end  = gdf[gdf['Academic Year'] == 'Endline']
                        genders_present = sorted(gdf['Gender'].dropna().astype(str).unique())
                        cols = st.columns(max(len(genders_present), 2))
                        for i, g in enumerate(genders_present):
                            with cols[i]:
                                b_mean = g_base[g_base['Gender'].astype(str) == g]['Obtained Marks'].mean() if not g_base.empty else None
                                e_mean = g_end[g_end['Gender'].astype(str) == g]['Obtained Marks'].mean()  if not g_end.empty  else None
                                if b_mean is not None and e_mean is not None:
                                    st.metric(f"{g} - Endline Avg", f"{e_mean:.2f}", delta=f"{e_mean - b_mean:.2f}")
                                elif e_mean is not None:
                                    st.metric(f"{g} - Endline Avg", f"{e_mean:.2f}")
                                elif b_mean is not None:
                                    st.metric(f"{g} - Baseline Avg", f"{b_mean:.2f}")

                        st.markdown("---")
                        gen_col1, gen_col2 = st.columns(2)
                        with gen_col1:
                            st.markdown("#### 📈 Average Score Trend")
                            st.caption("Direct comparison of mean scores by gender.")
                            avg_gen = gdf.groupby(['Gender', 'Academic Year'])['Obtained Marks'].mean().reset_index()
                            fig_gen_avg = px.bar(
                                avg_gen, x="Gender", y="Obtained Marks", color="Academic Year",
                                barmode="group", color_discrete_map=COLOR_MAP, text_auto='.2f',
                            )
                            fig_gen_avg.update_layout(
                                yaxis_title="Average Marks", margin=dict(l=0, r=0, t=30),
                                legend=dict(orientation="h", yanchor="top", y=-0.15, xanchor="center", x=0.5, title=""),
                            )
                            st.plotly_chart(fig_gen_avg, width='stretch')
                        with gen_col2:
                            st.markdown("#### 🧬 R.I.S.E Category Shift")
                            st.caption("Proportional breakdown of performance tiers by gender.")
                            gen_cat = gdf.groupby(['Gender', 'Academic Year', 'Category']).size().reset_index(name='Count')
                            gen_cat['Percentage'] = gen_cat.groupby(['Gender', 'Academic Year'])['Count'].transform(
                                lambda x: x / x.sum() * 100
                            )
                            fig_gen_rise = px.bar(
                                gen_cat, x="Academic Year", y="Percentage", color="Category", facet_col="Gender",
                                color_discrete_map=RISE_COLORS,
                                text=gen_cat['Percentage'].apply(lambda x: f'{x:.1f}%' if not pd.isna(x) and x > 5 else ''),
                                category_orders={
                                    "Category":      ["Reviving", "Initiating", "Shaping", "Evolving"],
                                    "Academic Year": ["Baseline", "Endline"],
                                },
                            )
                            fig_gen_rise.update_layout(
                                barmode='stack', yaxis_title="% of Students", margin=dict(l=0, r=0, t=40),
                                legend=dict(orientation="h", yanchor="top", y=-0.2, xanchor="center", x=0.5, title=""),
                            )
                            fig_gen_rise.update_xaxes(title_text='')
                            fig_gen_rise.for_each_annotation(lambda a: a.update(text=a.text.split("=")[-1]))
                            st.plotly_chart(fig_gen_rise, width='stretch')
                    else:
                        st.info("No valid gender data available in the current filtered selection.")
                else:
                    st.warning("⚠️ 'Gender' column is missing from the uploaded dataset.")

            # ------------------------------------------
            # TAB 6: RTM ANALYSIS
            # ------------------------------------------
            with tab6:
                st.markdown("### 📉 Regression to the Mean (RTM) Analysis")
                if not base_df.empty and not end_df.empty and 'Student ID' in df.columns:
                    base_rtm = base_df[['Student ID', 'Subject', 'Obtained Marks']].dropna(subset=['Student ID', 'Obtained Marks'])
                    end_rtm  = end_df[['Student ID',  'Subject', 'Obtained Marks']].dropna(subset=['Student ID', 'Obtained Marks'])
                    base_rtm = base_rtm.drop_duplicates(subset=['Student ID', 'Subject'])
                    end_rtm  = end_rtm.drop_duplicates(subset=['Student ID', 'Subject'])
                    rtm_df   = pd.merge(base_rtm, end_rtm, on=['Student ID', 'Subject'], suffixes=('_BL', '_EL'))

                    if not rtm_df.empty:
                        st.markdown("---")
                        normalize_rtm = st.checkbox(
                            "⚙️ Normalize scores (Z-scores) before analysis", value=False,
                            help="Standardizes scores so both have mean=0 and SD=1.",
                        )
                        if normalize_rtm:
                            rtm_df['Obtained Marks_BL'] = (rtm_df['Obtained Marks_BL'] - rtm_df['Obtained Marks_BL'].mean()) / rtm_df['Obtained Marks_BL'].std()
                            rtm_df['Obtained Marks_EL'] = (rtm_df['Obtained Marks_EL'] - rtm_df['Obtained Marks_EL'].mean()) / rtm_df['Obtained Marks_EL'].std()

                        rtm_df['Score Delta'] = rtm_df['Obtained Marks_EL'] - rtm_df['Obtained Marks_BL']
                        correlation = rtm_df['Obtained Marks_BL'].corr(rtm_df['Score Delta'])
                        variance    = rtm_df['Obtained Marks_BL'].var()
                        covariance  = rtm_df['Obtained Marks_BL'].cov(rtm_df['Score Delta'])
                        slope       = covariance / variance if variance and not pd.isna(variance) else 0.0
                        intercept   = rtm_df['Score Delta'].mean() - (slope * rtm_df['Obtained Marks_BL'].mean()) if not pd.isna(slope) else 0.0
                        total_rtm   = len(rtm_df)
                        improving_pct = len(rtm_df[rtm_df['Score Delta'] > 0]) / total_rtm * 100
                        declining_pct = len(rtm_df[rtm_df['Score Delta'] < 0]) / total_rtm * 100

                        if   slope <= -0.3: rtm_tag = "Strong RTM detected"
                        elif slope <= -0.1: rtm_tag = "Moderate RTM"
                        elif slope  <  0:   rtm_tag = "Minimal RTM"
                        else:               rtm_tag = "No RTM detected"

                        kpi_col1, kpi_col2, kpi_col3, kpi_col4 = st.columns(4)
                        kpi_col1.metric("Correlation (r)",        f"{correlation:.3f}" if not pd.isna(correlation) else "N/A")
                        kpi_col2.metric("Regression Slope",       f"{slope:.3f}")
                        kpi_col3.metric("Improving vs Declining", f"{improving_pct:.1f}% / {declining_pct:.1f}%")
                        kpi_col4.metric("Interpretation",         rtm_tag)

                        if slope <= -0.1:
                            st.warning("💡 **Important Insight:** Part of the observed improvement may be due to statistical regression to the mean rather than pure intervention impact.")
                        else:
                            st.success("💡 **Important Insight:** Observed improvements are less likely driven by RTM. The growth seen is more likely attributable to the actual impact of the educational intervention.")

                        st.markdown("---")
                        st.markdown("#### Core RTM View (Scatter Plot)")
                        st.caption("**Interpretation:** A **negative slope** (trendline going down) suggests RTM is present.")
                        fig_rtm = px.scatter(
                            rtm_df, x="Obtained Marks_BL", y="Score Delta",
                            trendline="ols", trendline_color_override="red", opacity=0.6,
                            color_discrete_sequence=["#636EFA"],
                            labels={
                                "Obtained Marks_BL": "Baseline Score (Z-Score)" if normalize_rtm else "Baseline Score",
                                "Score Delta": "Score Delta (Z-Score)" if normalize_rtm else "Score Delta (Endline - Baseline)",
                            },
                        )
                        fig_rtm.add_hline(
                            y=0, line_dash="dash", line_color="black",
                            annotation_text="No Change (Delta = 0)", annotation_position="bottom right",
                        )
                        fig_rtm.update_layout(margin=dict(l=0, r=0, t=30))
                        st.plotly_chart(fig_rtm, width='stretch')

                        st.markdown("---")
                        st.markdown("#### Binned Analysis (Quintiles)")
                        st.caption("Students are grouped into 5 equal-sized bins based on their initial Baseline scores.")
                        try:
                            rtm_df['BL_Quintile'] = pd.qcut(rtm_df['Obtained Marks_BL'], q=5, duplicates='drop')
                            binned_stats = rtm_df.groupby('BL_Quintile', observed=False).agg(
                                Avg_BL_Score=('Obtained Marks_BL', 'mean'),
                                Avg_Score_Delta=('Score Delta', 'mean'),
                                Student_Count=('Student ID', 'count'),
                            ).reset_index()
                            binned_stats['BL_Quintile_Str'] = binned_stats['BL_Quintile'].astype(str)
                            binned_stats = binned_stats.sort_values('Avg_BL_Score')
                            fig_binned = px.bar(
                                binned_stats, x='BL_Quintile_Str', y='Avg_Score_Delta',
                                text=binned_stats['Avg_Score_Delta'].apply(lambda x: f"{x:+.2f}"),
                                color='Avg_Score_Delta',
                                color_continuous_scale=px.colors.diverging.RdYlGn,
                                color_continuous_midpoint=0,
                                labels={
                                    "BL_Quintile_Str": "Baseline Score Range (Quintiles)",
                                    "Avg_Score_Delta": "Average Score Delta",
                                },
                                hover_data={"Student_Count": True, "Avg_BL_Score": ':.2f', "Avg_Score_Delta": ':.2f'},
                            )
                            fig_binned.add_hline(y=0, line_dash="solid", line_color="black", line_width=1)
                            fig_binned.update_traces(textposition='outside')
                            fig_binned.update_layout(margin=dict(l=0, r=0, t=30, b=40), coloraxis_showscale=False)
                            st.plotly_chart(fig_binned, width='stretch')
                        except ValueError:
                            st.info("Not enough variance in Baseline scores to generate quintile bins for this selection.")

                        st.markdown("---")
                        st.markdown("#### 🧮 Statistical Validation")
                        r_squared = correlation ** 2 if not pd.isna(correlation) else 0.0
                        sc1, sc2, sc3 = st.columns(3)
                        sc1.metric("Correlation (r)",      f"{correlation:.3f}" if not pd.isna(correlation) else "N/A")
                        sc2.metric("Regression Slope (b)", f"{slope:.3f}")
                        sc3.metric("R-squared (R²)",       f"{r_squared:.3f}")
                        st.markdown("**Mathematical Interpretation:**")
                        equation_str = f"**Score Delta = {intercept:.2f} + ({slope:.2f} × Baseline)**"
                        if slope < -0.3:
                            st.success(f"✔️ **Strong RTM Effect Confirmed:** {equation_str}")
                        elif slope < -0.1:
                            st.info(f"ℹ️ **Moderate RTM Effect:** {equation_str}")
                        elif slope < 0:
                            st.warning(f"⚠️ **Weak RTM Effect:** The slope is very close to zero. {equation_str}")
                        else:
                            st.error(f"❌ **No RTM Effect Detected:** The slope is positive ({slope:.2f}). {equation_str}")
                    else:
                        st.warning("⚠️ Could not find matching 'Student ID' and 'Subject' for RTM analysis.")
                else:
                    st.info("⚠️ Both Baseline and Endline datasets with a valid 'Student ID' column are required for this analysis.")

            # ==========================================
            # DRM REPORT GENERATION (PPTX)
            # ==========================================
            with st.sidebar:
                st.markdown("---")
                st.markdown("### 📄 DRM Compliance Report")
                if selected_donors != "All":
                    report_name = f"AY25-26_Impact_Report_{selected_donors.replace(' ', '_')}.pptx"

                    # FIX 6 — clear stale PPTX from session state when donor changes
                    if st.session_state.get('ready_ppt_donor') != selected_donors:
                        st.session_state.pop('ready_ppt', None)
                        st.session_state.pop('ready_ppt_donor', None)

                    if st.button(f"⚙️ Prepare PPTX for {selected_donors}", use_container_width=True):
                        with st.spinner("Compiling charts and generating presentation..."):
                            try:
                                from pptx import Presentation
                                from pptx.util import Inches, Pt

                                prs = Presentation()

                                def state_code(s):
                                    words = str(s).split()
                                    return "".join(w[0].upper() for w in words) if len(words) > 1 else str(s)[:2].upper()

                                def fig_to_png(fig):
                                    fig.update_layout(
                                        plot_bgcolor="white", paper_bgcolor="white", font=dict(size=13),
                                        yaxis=dict(
                                            showline=True, linecolor="black", linewidth=1.5,
                                            showticklabels=True, ticks="outside", tickcolor="black",
                                            ticklen=5, showgrid=True, gridcolor="lightgrey",
                                            gridwidth=1, zeroline=True, zerolinecolor="black",
                                            zerolinewidth=1, mirror=False,
                                        ),
                                        xaxis=dict(
                                            showline=True, linecolor="black", linewidth=1.5,
                                            showticklabels=True, ticks="outside", tickcolor="black",
                                            ticklen=5, showgrid=False, mirror=False,
                                        ),
                                    )
                                    buf = io.BytesIO()
                                    fig.write_image(buf, format="png", engine="kaleido", width=1000, height=550)
                                    buf.seek(0)
                                    return buf

                                # FIX 5 — track figures in an explicit dict instead of checking locals()
                                chart_figures = {}

                                def add_chart_slide(fig, title_text):
                                    slide = prs.slides.add_slide(prs.slide_layouts[5])
                                    slide.shapes.title.text = title_text
                                    slide.shapes.add_picture(fig_to_png(fig), Inches(0.5), Inches(1.5), width=Inches(9))

                                # Title slide
                                slide1 = prs.slides.add_slide(prs.slide_layouts[0])
                                slide1.shapes.title.text = "AY 25-26 Impact Report"
                                slide1.placeholders[1].text = (
                                    f"Donor: {selected_donors}\nGenerated automatically via Streamlit"
                                )

                                # Executive summary slide
                                slide2 = prs.slides.add_slide(prs.slide_layouts[1])
                                slide2.shapes.title.text = "Executive Summary"
                                tf = slide2.placeholders[1].text_frame
                                tf.word_wrap = True
                                num_schools       = filtered_df['Centre Name'].nunique()
                                subjects_assessed = ", ".join(sorted(filtered_df['Subject'].dropna().unique()))
                                states_in_data    = sorted(filtered_df['State'].dropna().unique())
                                states_str        = ", ".join(state_code(s) for s in states_in_data)
                                tf.text = f"States covered: {states_str}"
                                tf.add_paragraph().text = f"Total Centres Impacted: {num_schools}"
                                tf.add_paragraph().text = f"Subjects Assessed: {subjects_assessed}"
                                tf.add_paragraph().text = "Subject-wise Student Distribution (Endline):"
                                if not end_df.empty and 'Subject' in end_df.columns:
                                    for subj, count in (
                                        end_df.drop_duplicates(subset=['Student ID', 'Subject'])['Subject']
                                        .value_counts().sort_index().items()
                                    ):
                                        p = tf.add_paragraph()
                                        p.text  = f"{subj}: {count} students"
                                        p.level = 1
                                else:
                                    p = tf.add_paragraph()
                                    p.text  = "No endline data available."
                                    p.level = 1

                                # Build and register chart figures before adding slides
                                # Overall RISE shift
                                _cat_counts = filtered_df.groupby(['Academic Year', 'Category']).size().reset_index(name='Count')
                                _cat_counts['Percentage'] = _cat_counts.groupby('Academic Year')['Count'].transform(
                                    lambda x: x / x.sum() * 100
                                )
                                chart_figures['fig_rise'] = px.bar(
                                    _cat_counts, x="Category", y="Percentage", color="Academic Year",
                                    text=_cat_counts['Percentage'].apply(lambda x: f'{x:.1f}%' if not pd.isna(x) else ''),
                                    color_discrete_map=COLOR_MAP,
                                    category_orders={"Category": ["Reviving", "Initiating", "Shaping", "Evolving"]},
                                )
                                chart_figures['fig_rise'].update_layout(barmode='group')

                                # Box plot
                                chart_figures['fig_box'] = px.box(
                                    filtered_df, x="Academic Year", y="Obtained Marks",
                                    color="Academic Year", color_discrete_map=COLOR_MAP, points="all",
                                )
                                chart_figures['fig_box'].update_layout(showlegend=False)

                                # Gender avg
                                _gdf = filtered_df[
                                    ~filtered_df['Gender'].astype(str).str.lower().isin(['nan', 'none', 'null', ''])
                                ]
                                if not _gdf.empty:
                                    _avg_gen = _gdf.groupby(['Gender', 'Academic Year'])['Obtained Marks'].mean().reset_index()
                                    chart_figures['fig_gen_avg'] = px.bar(
                                        _avg_gen, x="Gender", y="Obtained Marks", color="Academic Year",
                                        barmode="group", color_discrete_map=COLOR_MAP, text_auto='.2f',
                                    )

                                # Add registered charts as slides
                                add_chart_slide(chart_figures['fig_rise'], "Overall R.I.S.E Category Shift (BL vs EL)")

                                # Per-subject grade RISE slides
                                all_subjects = sorted(filtered_df['Subject'].dropna().unique()) if 'Subject' in filtered_df.columns else []
                                for subj in all_subjects:
                                    subj_base = base_df[base_df['Subject'].astype(str) == subj] if 'Subject' in base_df.columns else pd.DataFrame()
                                    subj_end  = end_df[end_df['Subject'].astype(str) == subj]   if 'Subject' in end_df.columns  else pd.DataFrame()
                                    for period_df, label in [(subj_base, "Baseline"), (subj_end, "Endline")]:
                                        if not period_df.empty and 'Grade' in period_df.columns:
                                            grp = period_df.groupby(['Grade', 'Category']).size().reset_index(name='Count')
                                            grp['Percentage'] = grp.groupby('Grade')['Count'].transform(
                                                lambda x: x / x.sum() * 100
                                            )
                                            _fig = px.bar(
                                                grp, x="Grade", y="Percentage", color="Category",
                                                color_discrete_map=RISE_COLORS,
                                                text=grp['Percentage'].apply(lambda x: f'{x:.1f}%' if x > 5 else ''),
                                                category_orders={"Category": ["Reviving", "Initiating", "Shaping", "Evolving"]},
                                            )
                                            _fig.update_layout(
                                                barmode='stack', yaxis_title="% of Students",
                                                margin=dict(l=0, r=0, t=50), showlegend=True,
                                                legend=dict(orientation="h", yanchor="top", y=-0.2, xanchor="center", x=0.5, title=""),
                                            )
                                            add_chart_slide(_fig, f"{label} R.I.S.E by Grade — {subj}")

                                add_chart_slide(chart_figures['fig_box'], "Score Distribution (Box Plot)")
                                if 'fig_gen_avg' in chart_figures:
                                    add_chart_slide(chart_figures['fig_gen_avg'], "Average Score Trend by Gender")

                                ppt_buf = io.BytesIO()
                                prs.save(ppt_buf)
                                ppt_buf.seek(0)
                                st.session_state['ready_ppt']       = ppt_buf.getvalue()
                                st.session_state['ready_ppt_donor'] = selected_donors
                                st.success("Report prepared successfully!")

                            except ImportError:
                                st.error("⚠️ Missing required libraries! Please run: pip install python-pptx kaleido")
                            except Exception as e:
                                st.error(f"⚠️ Error preparing presentation: {e}")

                    if 'ready_ppt' in st.session_state and st.session_state.get('ready_ppt_donor') == selected_donors:
                        st.download_button(
                            label="⬇️ Download Presentation",
                            data=st.session_state['ready_ppt'],
                            file_name=report_name,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                        )
                else:
                    st.info("💡 Select a specific Donor from the global filters to enable the DRM Report generator.")

else:
    st.warning(
        "⚠️ No data file loaded. Please upload an Excel file above to populate the dashboard."
    )
